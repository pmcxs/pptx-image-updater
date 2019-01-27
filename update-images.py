from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.action import PP_ACTION

import ssl
import base64
import argparse

try:
    from urllib.request import Request, urlopen  # Python 3
except ImportError:
    from urllib2 import Request, urlopen  # Python 2

parser = argparse.ArgumentParser()
parser.add_argument("pptx_path", type=str, help="input path to the pptx file that will be processed")
parser.add_argument("-u","--username", help="username that will be used to send basic auth credentials")
parser.add_argument("-p","--password", help="password that will be used to send basic auth credentials. Please enclose with single quotes if it has special characters")
args = parser.parse_args()

try:
    prs = Presentation(args.pptx_path)
except Exception as e:
    print("Couldn't open specified file. Error: %s" % str(e))
    exit()

if args.username and args.password:
    print("Setting basic authorization")
    authorizationHeaderValue = "Basic " + (base64.b64encode(str(args.username + ":" + args.password).encode('utf-8'))).decode('utf-8')
else:
    authorizationHeaderValue = ""

for sx, slide in enumerate(prs.slides):
    print("------------------\nProcessing slide %d\n------------------" % (sx+1))

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

            x, y, cx, cy = shape.left, shape.top, shape.width, shape.height

            print("Found picture at coordinates: (%d:%d)" % (x,y))
         
            click_action = shape.click_action

            if click_action.action == PP_ACTION.HYPERLINK:

                url = click_action.hyperlink.address
                print("\tReplacing with image from: %s" % url)

                try:
                    q = Request(url)

                    if authorizationHeaderValue:
                        q.add_header('Authorization', authorizationHeaderValue)

                    new_picture_blob = urlopen(q, context=ssl._create_unverified_context()).read()

                    imgPic = shape._pic
                    imgRID = imgPic.blipFill.blip.attrib.values()[0] 
                    imgPart = slide.part.related_parts[imgRID] 
                    imgPart._blob = new_picture_blob

                except Exception as e:
                    print("\tCouldn't load and update image. Error: %s" % str(e))
            else:
                print("\tImage doesn't have an hyperlink. Ignoring it")
try:
    prs.save(args.pptx_path)
except Exception as e:
    print("Couldn't save powerpoint file: " + str(e))